# -------------------------------------------------------------
# Streamlit App: PPT → Voice Preview → Download PPT with Voice
# -------------------------------------------------------------

# ===================== IMPORTS =====================
import os
import tempfile
import time
from pathlib import Path

import streamlit as st
from pptx import Presentation
from pptx.util import Inches

from dotenv import load_dotenv
from openai import OpenAI, RateLimitError, APIStatusError
from pydub import AudioSegment

# ===================== ENV ========================
load_dotenv()
OPENAI_API_KEY = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")
ANTHROPIC_API_KEY = st.secrets.get("ANTHROPIC_API_KEY") or os.getenv("ANTHROPIC_API_KEY")

if not OPENAI_API_KEY:
    st.error("❌ OPENAI_API_KEY not configured")
    st.stop()

client = OpenAI(api_key=OPENAI_API_KEY)

# Claude is only used as a fallback for narration TEXT when OpenAI's quota or
# rate limit is hit — Anthropic has no text-to-speech API, so it can't help
# with the voice-over audio itself (that fallback uses gTTS, see below).
anthropic_client = None
CLAUDE_MODEL = "claude-haiku-4-5-20251001"
if ANTHROPIC_API_KEY:
    import anthropic

    anthropic_client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# ================= UI SETUP ======================
st.set_page_config(page_title="PPT Voice Over Studio", layout="wide")
st.title("🎤 PPT Voice Over Studio")
st.caption("Title + Content based narration • Voice & Pitch Control")

st.divider()

# ================= SESSION STATE =================
if "slides" not in st.session_state:
    st.session_state.slides = []
if "ppt_loaded" not in st.session_state:
    st.session_state.ppt_loaded = False
if "ppt_path" not in st.session_state:
    st.session_state.ppt_path = None
if "ppt_name" not in st.session_state:
    st.session_state.ppt_name = None

# ================= SIDEBAR CONTROLS =================
st.sidebar.header("🎙 Voice Settings")

voice_choice = st.sidebar.selectbox("Select Voice", ["Female", "Male"])

pitch = st.sidebar.slider(
    "Voice Pitch",
    min_value=-6,
    max_value=6,
    value=0,
    help="Negative = deeper voice, Positive = sharper voice",
)

if not anthropic_client:
    st.sidebar.caption(
        "ℹ️ Add ANTHROPIC_API_KEY to enable a Claude fallback for narration "
        "text if OpenAI's quota or rate limit is reached."
    )

VOICE_MAP = {
    "Male": "alloy",
    "Female": "verse",
}

# ================= HELPERS =======================
def get_slide_title(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    return "this concept"


# 🔑 ONLY CHANGE IS HERE (dynamic openings)
OPENING_TEMPLATES = [
    "Today we are going to explore {title}. ",
    "Let us now understand {title}. ",
    "In this section, we will focus on {title}. ",
    "Next, we are going to look at {title}. ",
    "Here, we will discuss {title}. ",
]


def _is_openai_limit_error(exc: Exception) -> bool:
    """True when OpenAI rejected the call for quota / rate-limit reasons."""
    if isinstance(exc, RateLimitError):
        return True
    if isinstance(exc, APIStatusError) and exc.status_code == 429:
        return True
    msg = str(exc).lower()
    return "insufficient_quota" in msg or "rate_limit" in msg or "quota" in msg


def generate_narration(slide_text: str, slide_index: int, slide_title: str) -> str:
    title = slide_title.strip()

    if slide_index == 0:
        opening = f"Today we are going to explore {title}. "
    else:
        template = OPENING_TEMPLATES[slide_index % len(OPENING_TEMPLATES)]
        opening = template.format(title=title)

    prompt = f"""
You are narrating a PowerPoint slide.

STRICT RULES:
- Use the slide title EXACTLY as given
- NEVER say "this slide", "the topic", or generic phrases
- Speak ONLY about the slide title and slide content
- Simple Indian teaching tone
- No headings
- No bullet points

Start exactly with:
"{opening}"

Slide Title:
{title}

Slide Content:
{slide_text}
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
        )
        return response.choices[0].message.content.strip()
    except Exception as exc:
        if not _is_openai_limit_error(exc):
            raise
        return _generate_narration_claude(prompt)


def _generate_narration_claude(prompt: str) -> str:
    if not anthropic_client:
        raise RuntimeError(
            "OpenAI's quota/rate limit was reached and no ANTHROPIC_API_KEY is "
            "configured. Add one to .env (or .streamlit/secrets.toml) to enable "
            "the Claude fallback."
        )
    st.warning("⚠️ OpenAI limit reached — using Claude to write this narration.")
    message = anthropic_client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=600,
        messages=[{"role": "user", "content": prompt}],
    )
    return "".join(
        block.text for block in message.content if block.type == "text"
    ).strip()

# ================= SAFE TTS ======================
def chunk_text(text, max_chars=900):
    chunks, current = [], ""
    for sentence in text.split(". "):
        if len(current) + len(sentence) < max_chars:
            current += sentence + ". "
        else:
            chunks.append(current.strip())
            current = sentence + ". "
    if current.strip():
        chunks.append(current.strip())
    return chunks


def apply_pitch(audio_path: Path, pitch_change: int):
    if pitch_change == 0:
        return audio_path

    audio = AudioSegment.from_mp3(audio_path)
    new_sample_rate = int(audio.frame_rate * (2.0 ** (pitch_change / 12.0)))
    pitched = audio._spawn(audio.raw_data, overrides={"frame_rate": new_sample_rate})
    pitched = pitched.set_frame_rate(44100)
    pitched.export(audio_path, format="mp3")
    return audio_path


def _gtts_fallback(text: str, out_mp3: Path):
    """Free, keyless TTS used when OpenAI's voice quota/rate limit is hit.

    gTTS has a single voice per language, so voice choice (male/female) has
    no effect here — pitch shifting is still applied afterwards.
    """
    from gtts import gTTS

    combined = AudioSegment.empty()
    for chunk in chunk_text(text):
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
            tmp_path = Path(tmp.name)
        try:
            gTTS(text=chunk, lang="en").save(tmp_path)
            combined += AudioSegment.from_mp3(tmp_path)
        finally:
            tmp_path.unlink(missing_ok=True)
    combined.export(out_mp3, format="mp3")


def openai_tts(text: str, out_mp3: Path, voice: str, pitch_change: int, retries=3):
    chunks = chunk_text(text)

    try:
        with open(out_mp3, "wb") as f:
            for chunk in chunks:
                succeeded = False
                for _ in range(retries):
                    try:
                        with client.audio.speech.with_streaming_response.create(
                            model="gpt-4o-mini-tts",
                            voice=voice,
                            input=chunk,
                        ) as response:
                            for audio_bytes in response.iter_bytes():
                                f.write(audio_bytes)
                        succeeded = True
                        break
                    except Exception as exc:
                        if _is_openai_limit_error(exc):
                            raise
                        time.sleep(1)
                if not succeeded:
                    raise RuntimeError("OpenAI TTS failed after retries")
    except Exception as exc:
        if not _is_openai_limit_error(exc):
            raise
        st.warning(
            "⚠️ OpenAI voice limit reached — using free Google TTS instead "
            "(voice choice doesn't apply, pitch styling still does)."
        )
        _gtts_fallback(text, out_mp3)
        apply_pitch(out_mp3, pitch_change)
        return

    apply_pitch(out_mp3, pitch_change)


def add_audio_to_slide(slide, audio_path: Path):
    slide.shapes.add_movie(
        movie_file=str(audio_path),
        left=Inches(0.3),
        top=Inches(0.3),
        width=Inches(1),
        height=Inches(1),
        mime_type="audio/mpeg",
    )

# ================= FILE UPLOAD ====================
ppt_file = st.file_uploader("📤 Upload PPTX", type=["pptx"])

if ppt_file and not st.session_state.ppt_loaded:
    workdir = Path(tempfile.mkdtemp())
    ppt_path = workdir / ppt_file.name
    ppt_path.write_bytes(ppt_file.read())

    prs = Presentation(ppt_path)
    st.session_state.slides.clear()

    for idx, slide in enumerate(prs.slides):
        slide_text = " ".join(
            shape.text
            for shape in slide.shapes
            if hasattr(shape, "text") and shape != slide.shapes.title
        ).strip()

        slide_title = get_slide_title(slide)

        notes = generate_narration(slide_text, idx, slide_title)

        st.session_state.slides.append({
            "index": idx,
            "text": slide_text or slide_title,
            "notes": notes,
            "skip": False,
        })

    st.session_state.ppt_loaded = True
    st.session_state.ppt_path = ppt_path
    st.session_state.ppt_name = ppt_file.name
    st.success("✅ PPT loaded successfully")

# ================= PREVIEW ========================
if st.session_state.ppt_loaded:
    st.subheader("🎧 Preview Voice")

    for slide in st.session_state.slides:
        with st.expander(f"Slide {slide['index'] + 1}"):
            slide["notes"] = st.text_area(
                "Narration Text",
                slide["notes"],
                key=f"notes_{slide['index']}",
                height=130,
            )

            if st.button("▶ Preview Voice", key=f"preview_{slide['index']}"):
                with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as f:
                    openai_tts(
                        slide["notes"],
                        Path(f.name),
                        VOICE_MAP[voice_choice],
                        pitch,
                    )
                    st.audio(f.name)

# ================= FINAL GENERATION =================
st.divider()

if st.session_state.ppt_loaded:
    if st.button("📥 Download PPT with Voice-over"):
        prs = Presentation(st.session_state.ppt_path)
        outdir = Path(tempfile.mkdtemp())

        total = len(st.session_state.slides)
        progress = st.progress(0.0)

        for i, slide_data in enumerate(st.session_state.slides, start=1):
            progress.progress(i / total)

            slide = prs.slides[slide_data["index"]]
            mp3_path = outdir / f"slide_{slide_data['index']}.mp3"

            openai_tts(
                slide_data["notes"],
                mp3_path,
                VOICE_MAP[voice_choice],
                pitch,
            )

            add_audio_to_slide(slide, mp3_path)

            try:
                slide.notes_slide.placeholders[1].text = slide_data["notes"]
            except Exception:
                pass

        final_ppt = outdir / st.session_state.ppt_name
        prs.save(final_ppt)

        st.download_button(
            "⬇ Download PPT with Voice-over",
            final_ppt.read_bytes(),
            file_name=st.session_state.ppt_name,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
